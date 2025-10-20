"""PPTX metadata extractor using python-pptx."""

from typing import ClassVar

from pptx import Presentation

from metaextract.core.models import DocumentMetadata, ExtractionResult, FileType
from metaextract.extractors.base import MetadataExtractor


class PPTXExtractor(MetadataExtractor):
    """Extract metadata from PPTX files using python-pptx."""

    SUPPORTED_TYPES: ClassVar[list[FileType]] = [FileType.PPTX]

    def extract(self) -> ExtractionResult:
        """Extract metadata from PPTX."""
        try:
            metadata = self._create_base_metadata()

            prs = Presentation(str(self.file_path))
            core_props = prs.core_properties

            # Extract author
            if core_props.author:
                metadata.author = self._clean_string(core_props.author)
                if metadata.author and metadata.author not in metadata.users:
                    metadata.users.append(metadata.author)

            # Extract last modified by
            if core_props.last_modified_by:
                metadata.last_modified_by = self._clean_string(core_props.last_modified_by)
                if metadata.last_modified_by and metadata.last_modified_by not in metadata.users:
                    metadata.users.append(metadata.last_modified_by)

            # Extract dates
            metadata.created = core_props.created
            metadata.modified = core_props.modified

            # Extract application info
            self._extract_app_info(metadata)

            # Store raw properties
            metadata.raw = {
                "author": str(core_props.author) if core_props.author else None,
                "last_modified_by": str(core_props.last_modified_by) if core_props.last_modified_by else None,
                "created": str(core_props.created) if core_props.created else None,
                "modified": str(core_props.modified) if core_props.modified else None,
                "title": str(core_props.title) if core_props.title else None,
                "subject": str(core_props.subject) if core_props.subject else None,
                "category": str(core_props.category) if core_props.category else None,
                "revision": str(core_props.revision) if core_props.revision else None,
            }

            self._metadata = metadata
            return ExtractionResult(success=True, metadata=metadata)

        except Exception as e:
            return ExtractionResult(success=False, error=str(e))

    def _extract_app_info(self, metadata: DocumentMetadata) -> None:
        """Extract application info from the PPTX package."""
        import zipfile

        from lxml import etree

        try:
            with zipfile.ZipFile(self.file_path, "r") as zf:
                if "docProps/app.xml" in zf.namelist():
                    app_xml = zf.read("docProps/app.xml")
                    root = etree.fromstring(app_xml)

                    ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}

                    app_elem = root.find("ep:Application", ns)
                    if app_elem is not None and app_elem.text:
                        metadata.application = app_elem.text
                        if app_elem.text not in metadata.software:
                            metadata.software.append(app_elem.text)

                    version_elem = root.find("ep:AppVersion", ns)
                    if version_elem is not None and version_elem.text:
                        metadata.app_version = version_elem.text

                    # Check for template
                    template_elem = root.find("ep:Template", ns)
                    if template_elem is not None and template_elem.text:
                        metadata.template = template_elem.text

        except Exception:
            pass

    def extract_text(self) -> str | None:
        """Extract text content from slides."""
        try:
            prs = Presentation(str(self.file_path))
            texts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)

            return "\n".join(texts)
        except Exception:
            return None
